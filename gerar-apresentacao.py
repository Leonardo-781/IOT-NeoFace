#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para gerar apresentação PowerPoint da Etapa 1
Painel IoT Distribuído com Monitoramento em Tempo Real
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Cores do tema
PRIMARY = RGBColor(18, 56, 102)      # Azul escuro
SECONDARY = RGBColor(28, 120, 86)    # Verde
LIGHT_BLUE = RGBColor(236, 242, 250)
LIGHT_GREEN = RGBColor(235, 247, 241)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle, author, date_str):
    """Adiciona slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Autor
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.5))
    author_frame = author_box.text_frame
    author_p = author_frame.paragraphs[0]
    author_p.text = author
    author_p.font.size = Pt(24)
    author_p.font.color.rgb = WHITE
    author_p.alignment = PP_ALIGN.CENTER
    
    # Data
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = date_str
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = RGBColor(200, 200, 200)
    date_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_func):
    """Adiciona slide de conteúdo com background branco"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título com barra verde
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY
    title_shape.line.color.rgb = PRIMARY
    
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    title_frame.margin_top = Inches(0.1)
    
    # Chamar função de conteúdo
    content_func(slide)

def slide_micromundo(slide):
    y = 1.2
    
    # Left column - Características
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(4.3), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_BLUE
    left_box.line.color.rgb = SECONDARY
    left_box.line.width = Pt(3)
    
    text_frame = left_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Características"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY
    
    items = [
        "Coleta contínua de dados",
        "Múltiplos nós geograficamente distribuídos",
        "Sincronização centralizada",
        "Dashboard web",
        "Displays locais OLED"
    ]
    for item in items:
        p = text_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
    
    # Right column - Sensores
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(y), Inches(4.3), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GREEN
    right_box.line.color.rgb = SECONDARY
    right_box.line.width = Pt(3)
    
    text_frame = right_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Sensores"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY
    
    sensors = ["Temperatura", "Umidade", "Pressão", "Força de sinal WiFi (RSSI)", "Qualidade de conexão"]
    for sensor in sensors:
        p = text_frame.add_paragraph()
        p.text = "• " + sensor
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)

def slide_o_que_faz(slide):
    y = 1.2
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(5.8))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "1. Coleta contínua de leituras de sensores (temperatura, umidade, pressão)",
        "2. Monitoramento de força e qualidade de sinal WiFi para diagnóstico",
        "3. Exibição em tempo real em telas OLED com animações responsivas",
        "4. Controle interativo via sensor de toque capacitivo",
        "5. Sincronização com servidor central através de API HTTP",
        "6. Persistência de séries temporais em PostgreSQL",
        "7. Dashboard web para visualização e controle remoto"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)

def slide_arquitetura(slide):
    y = 1.2
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(5.8))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "ESP32/ESP8266 → [WiFi] → Nginx → API Node.js → PostgreSQL"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(30)
    
    # Diagrama simplificado
    p = text_frame.add_paragraph()
    p.text = ""
    
    diagram_items = [
        ("ESP32 (Sala)", "WiFi POST"),
        ("Nginx", "HTTP"),
        ("API Node.js", "TCP"),
        ("PostgreSQL", "")
    ]
    
    for item, protocol in diagram_items:
        p = text_frame.add_paragraph()
        if protocol:
            p.text = f"  {item}  ─{protocol}→"
        else:
            p.text = f"  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY
        p.space_before = Pt(8)

def slide_infraestrutura(slide):
    y = 1.2
    
    # RPi
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(3.0), Inches(2.0))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.color.rgb = PRIMARY
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Raspberry Pi 3B"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    
    for item in ["PostgreSQL", "DB persistente", "IP: 192.168.1.10", "Ethernet"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
    
    # Notebook
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(y), Inches(3.0), Inches(2.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GREEN
    box2.line.color.rgb = SECONDARY
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Notebook (Linux Mint)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    
    for item in ["Node.js API", "Nginx proxy", "IP: 192.168.1.11", "WiFi"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
    
    # ESP nodes
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(y), Inches(2.5), Inches(2.0))
    box3.fill.solid()
    box3.fill.fore_color.rgb = LIGHT_BLUE
    box3.line.color.rgb = PRIMARY
    box3.line.width = Pt(2)
    
    tf = box3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ESP32/ESP8266"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    
    for item in ["Múltiplos nós", "Sensores", "WiFi", "OLED local"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
    
    # Fluxo
    y = 3.5
    text_box = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(9.2), Inches(3.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Protocolo: HTTP POST com payload JSON"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)
    
    features = [
        ("Frequência", "30-60s"),
        ("Buffer local", "10 leituras"),
        ("Auto-recover", "WiFi com backoff"),
        ("Rate limiting", "CORS + compressão gzip")
    ]
    
    for feat, desc in features:
        p = text_frame.add_paragraph()
        p.text = f"• {feat}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)

def slide_sensores(slide):
    y = 1.2
    
    # BMP280
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(4.5), Inches(2.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.color.rgb = PRIMARY
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BMP280"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    
    for item in ["Temp: -40 a +85°C", "Pressão: 300-1100 hPa", "Proto: I2C/SPI", "Precisão: ±1°C"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
    
    # DHT11
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(y), Inches(4.5), Inches(2.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GREEN
    box2.line.color.rgb = SECONDARY
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DHT11"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    
    for item in ["Temp: 0-50°C", "Umidade: 20-80%", "Proto: Single-wire", "Precisão: ±2°C, ±5%"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
    
    # RSSI
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.3), Inches(9.2), Inches(2.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = LIGHT_BLUE
    box3.line.color.rgb = PRIMARY
    box3.line.width = Pt(2)
    
    tf = box3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "RSSI (Received Signal Strength Indicator)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    
    p = tf.add_paragraph()
    p.text = "Monitoramento de qualidade de sinal WiFi:"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(8)
    
    for item in ["• Diagnóstico real de cobertura", "• Detecção de interferência", "• Planejamento de repetidores"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

def slide_displays(slide):
    y = 1.2
    
    # OLED
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(6.0), Inches(5.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.color.rgb = PRIMARY
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Display OLED SSD1306"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    
    for item in ["128×64 pixels", "Protocolo I2C", "Taxa: 1-2 Hz", "Temperatura real", "Força WiFi", "Animações rosto", "Status sistema"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
    
    # TTP223
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(y), Inches(3.0), Inches(5.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GREEN
    box2.line.color.rgb = SECONDARY
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sensor TTP223"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    
    for item in ["Capacitivo", "5 expressões", "Debounce 300ms", "Ciclagem ao tocar"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "😊 😢 😲"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

def slide_postgresql(slide):
    y = 1.2
    
    # Telemetry
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(9.2), Inches(1.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.color.rgb = PRIMARY
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tabela: telemetry (TimescaleDB)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    
    for item in ["time (TIMESTAMP), device_id, temperature_c, humidity_percent, pressure_hpa", "rssi_dbm, battery_voltage, Retenção: 90 dias"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
    
    # Devices
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.2), Inches(9.2), Inches(1.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GREEN
    box2.line.color.rgb = SECONDARY
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tabela: devices"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    
    p = tf.add_paragraph()
    p.text = "• ID, nome, localização, status, last_seen, firmware_version, ip_address"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    
    # Alerts
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.0), Inches(9.2), Inches(1.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = LIGHT_BLUE
    box3.line.color.rgb = PRIMARY
    box3.line.width = Pt(2)
    
    tf = box3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tabela: alerts"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    
    p = tf.add_paragraph()
    p.text = "• ID, device_id, alert_type, severity, created_at, resolved_at"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

def slide_objetivos(slide):
    y = 1.2
    
    objectives = [
        ("Escalabilidade", "Novos nós podem ser adicionados sem redesenho. Nginx permite escalabilidade horizontal."),
        ("Disponibilidade", "Buffer local em ESP, Nginx e API redundantes, PostgreSQL com snapshots automáticos."),
        ("Tolerância a Falhas", "Reconexão automática, retry com backoff, health check, fila de mensagens locais."),
        ("Compartilhamento", "Nginx como ponto centralizado, PostgreSQL como estado único, transações ACID.")
    ]
    
    for i, (title, desc) in enumerate(objectives):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y + i * 1.35), Inches(9.2), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else LIGHT_GREEN
        box.line.color.rgb = PRIMARY if i % 2 == 0 else SECONDARY
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {title}"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = PRIMARY
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

def slide_desafios(slide):
    y = 1.2
    
    # Desafios
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(4.5), Inches(5.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    box1.line.color.rgb = PRIMARY
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Infraestrutura"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)
    
    for item in ["IP dinâmico", "Falhas WiFi", "Latência variável", "Nós dispersos"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Soluções
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(y), Inches(4.5), Inches(5.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_GREEN
    box2.line.color.rgb = SECONDARY
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Soluções"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)
    
    for item in ["DHCP estático", "Buffer local", "Backoff exponencial", "Sincronização periódica"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

def slide_roadmap(slide):
    y = 1.2
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(5.8))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    stages = [
        ("Etapa 2", "Threads em Node.js, processamento paralelo"),
        ("Etapa 3", "UDP para comunicação de baixa latência"),
        ("Etapa 4", "RPC entre servidores"),
        ("Etapa 5", "Message Bus (RabbitMQ/MQTT)"),
        ("Etapa 6", "Service Discovery (mDNS/Consul)"),
        ("Etapa 7", "Mutual exclusion, eleição de líder"),
        ("Etapa 8", "Replicação com consistência eventual"),
        ("Etapa 9", "Circuit breaker, resiliência"),
        ("Etapa 10", "TLS, autenticação mútua, JWT")
    ]
    
    for i, (stage, desc) in enumerate(stages):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"  {stage}: {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(3)

def slide_conclusao(slide):
    y = 1.5
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = PRIMARY
    box.line.width = Pt(3)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Um Sistema Distribuído Real"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)
    
    items = [
        "Infraestrutura genuinamente prática e doméstica",
        "Múltiplos nós heterogêneos (Raspberry, Notebook, ESP32/ESP8266)",
        "Rede WiFi com falhas esperadas",
        "Persistência centralizada",
        "Desafios reais de sincronização e resiliência"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "▸ " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1.0))
    footer_frame = footer_box.text_frame
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Fundação Sólida para Conceitos Avançados"
    footer_p.font.size = Pt(18)
    footer_p.font.bold = True
    footer_p.font.color.rgb = SECONDARY
    footer_p.alignment = PP_ALIGN.CENTER

# Criar slides
print("Criando apresentação PPTX...")

# Slide 1: Título
add_title_slide(prs, 
    "Painel IoT Distribuído",
    "com Monitoramento em Tempo Real",
    "Leonardo Cardoso de Moura | Matrícula: 32021BSI004",
    "08 de maio de 2026")

# Slide 2: O Micro-mundo
add_content_slide(prs, "O Micro-mundo", slide_micromundo)

# Slide 3: O que faz
add_content_slide(prs, "O que o Sistema Faz", slide_o_que_faz)

# Slide 4: Arquitetura
add_content_slide(prs, "Arquitetura Geral", slide_arquitetura)

# Slide 5: Infraestrutura
add_content_slide(prs, "Infraestrutura Física", slide_infraestrutura)

# Slide 6: Sensores
add_content_slide(prs, "Sensores Instalados", slide_sensores)

# Slide 7: Displays
add_content_slide(prs, "Displays e Interatividade", slide_displays)

# Slide 8: PostgreSQL
add_content_slide(prs, "Schema PostgreSQL", slide_postgresql)

# Slide 9: Objetivos
add_content_slide(prs, "Objetivos de Sistemas Distribuídos", slide_objetivos)

# Slide 10: Desafios
add_content_slide(prs, "Desafios Reais do Sistema", slide_desafios)

# Slide 11: Roadmap
add_content_slide(prs, "Próximas Etapas de Desenvolvimento", slide_roadmap)

# Slide 12: Conclusão
add_content_slide(prs, "Conclusão", slide_conclusao)

# Slide 13: Obrigado
add_title_slide(prs,
    "Obrigado!",
    "",
    "Leonardo Cardoso de Moura",
    "Disciplina: Sistemas Distribuídos | 08 de maio de 2026")

# Salvar
output_path = r"c:\Users\Leonardo\OneDrive\Documentos\VS Code\Sistemas Distribuidos\Trabalho Final\etapa-1\apresentacao\apresentacao.pptx"
prs.save(output_path)
print(f"✓ Apresentação criada com sucesso!")
print(f"✓ Salvo em: {output_path}")
